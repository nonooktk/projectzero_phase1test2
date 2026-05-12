"""
バックエンド処理説明スライド生成スクリプト（2スライド構成）
  Slide1: ChromaDB × GraphRAG の役割と組み合わせ効果
  Slide2: AI分析への接続と可視化
実行: python docs/generate_backend_slide.py
出力: docs/backend_slide.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUTPUT_PATH = "docs/backend_slide.pptx"

C_BG     = RGBColor(0xFF, 0xFF, 0xFF)
C_NAVY   = RGBColor(0x1B, 0x2E, 0x45)
C_TEAL   = RGBColor(0x00, 0x99, 0xCC)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_PURPLE = RGBColor(0x7C, 0x3A, 0xED)   # ChromaDB
C_RUST   = RGBColor(0xC0, 0x50, 0x10)   # GraphRAG
C_GREEN  = RGBColor(0x16, 0xA3, 0x4A)
C_AMBER  = RGBColor(0xD9, 0x77, 0x06)
C_GRAY   = RGBColor(0x9C, 0xA3, 0xAF)
C_LIGHT  = RGBColor(0xF3, 0xF4, 0xF6)
C_BORDER = RGBColor(0xD1, 0xD5, 0xDB)
C_PALE_P = RGBColor(0xED, 0xE9, 0xFE)   # 薄紫（ChromaDB 背景）
C_PALE_R = RGBColor(0xFF, 0xED, 0xCC)   # 薄橙（GraphRAG 背景）
C_SUB    = RGBColor(0xCC, 0xDD, 0xFF)   # サブテキスト（dark box内）

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
SW = SLIDE_W - Inches(0.30)

HDR_Y, HDR_H = Inches(0.10), Inches(0.42)
FTR_H = Inches(0.40)
FTR_Y = SLIDE_H - FTR_H


# ── ユーティリティ ──────────────────────────────

def rect(slide, x, y, w, h, fill=None, border=None, bw=Pt(1.2), dashed=False):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
        if dashed:
            ln = s.line._ln
            for c in ln.findall(qn('a:prstDash')): ln.remove(c)
            pd = etree.SubElement(ln, qn('a:prstDash')); pd.set('val', 'dash')
    else:
        s.line.fill.background()
    return s


def t(slide, text, x, y, w, h, size=Pt(11), color=None,
      bold=False, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold
    r.font.color.rgb = color or C_NAVY
    return tb


def arrowhead(conn):
    ln = conn.line._ln
    el = ln.find(qn('a:tailEnd'))
    if el is None:
        el = etree.SubElement(ln, qn('a:tailEnd'))
    el.set('type', 'arrow'); el.set('w', 'med'); el.set('len', 'med')


def harrow(slide, x1, y, x2, color=None, lbl=None, above=True):
    c = color or C_NAVY
    co = slide.shapes.add_connector(1, x1, y, x2, y)
    co.line.color.rgb = c; co.line.width = Pt(1.5)
    arrowhead(co)
    if lbl:
        lx, lw = min(x1, x2), abs(x2 - x1)
        ly = y - Inches(0.20) if above else y + Inches(0.03)
        t(slide, lbl, lx, ly, lw, Inches(0.18), size=Pt(8), color=c, bold=True)


def varrow(slide, x, y1, y2, color=None):
    c = color or C_NAVY
    co = slide.shapes.add_connector(1, x, y1, x, y2)
    co.line.color.rgb = c; co.line.width = Pt(1.5)
    arrowhead(co)


def chip(slide, label, x, y, w, h, fill=None, sub=None, sub_color=None):
    """塗り潰しボックス（サブテキスト付き対応）"""
    bg = fill or C_NAVY
    rect(slide, x, y, w, h, fill=bg)
    if sub:
        t(slide, label, x + Inches(0.08), y + Inches(0.04), w - Inches(0.16), h * 0.45,
          size=Pt(10), color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)
        t(slide, sub, x + Inches(0.08), y + h * 0.46, w - Inches(0.16), h * 0.50,
          size=Pt(9), color=sub_color or C_SUB, align=PP_ALIGN.LEFT)
    else:
        t(slide, label, x, y, w, h, size=Pt(10), color=C_WHITE, bold=True)


def panel(slide, title, x, y, w, h, hdr_color):
    """パネル枠（ヘッダー付き）を描画し、コンテンツ開始Y を返す"""
    rect(slide, x, y, w, h, fill=C_LIGHT, border=C_BORDER)
    rect(slide, x, y, w, Inches(0.36), fill=hdr_color)
    t(slide, title, x + Inches(0.12), y, w - Inches(0.20), Inches(0.36),
      size=Pt(12), color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)
    return y + Inches(0.44)   # コンテンツ開始Y


def section_label(slide, text, x, y, w, color):
    """セクション内の小見出し"""
    t(slide, text, x, y, w, Inches(0.24),
      size=Pt(9), color=color, bold=True, align=PP_ALIGN.LEFT)


def hdr(slide, msg):
    rect(slide, Inches(0.15), HDR_Y, SW, HDR_H, fill=C_NAVY)
    t(slide, msg, Inches(0.15), HDR_Y, SW, HDR_H,
      size=Pt(13), color=C_WHITE, bold=True)


def ftr(slide, msg):
    rect(slide, Inches(0.15), FTR_Y, SW, FTR_H, fill=C_NAVY)
    t(slide, msg, Inches(0.15), FTR_Y, SW, FTR_H,
      size=Pt(10), color=C_WHITE)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ══════════════════════════════════════════════════════════════════════
# Slide 1 : ChromaDB × GraphRAG の役割と組み合わせ効果
# ══════════════════════════════════════════════════════════════════════

def build_slide1(prs):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG)

    hdr(slide, "ChromaDB × GraphRAG の併用で「関連度」と「関連性」を同時に捉え、抜け漏れのない情報収集を実現する")

    # ── レイアウト定数 ──
    panel_y   = HDR_Y + HDR_H + Inches(0.12)
    panel_h   = Inches(4.50)
    lx, lw    = Inches(0.15), Inches(5.65)     # 左パネル（ChromaDB）
    mx, mw    = Inches(5.88), Inches(1.57)     # 中央（組み合わせ効果）
    rx, rw    = Inches(7.53), Inches(5.65)     # 右パネル（GraphRAG）
    tbl_y     = panel_y + panel_h + Inches(0.12)
    tbl_h     = FTR_Y - tbl_y - Inches(0.08)

    # ════════════════════════════════
    # 左パネル：ChromaDB（関連度）
    # ════════════════════════════════
    cy_l = panel(slide, "ChromaDB ── 関連度（意味の近さ）で検索", lx, panel_y, lw, panel_h, C_PURPLE)
    pad = Inches(0.12)
    iw = lw - pad * 2

    # 仕組み
    section_label(slide, "▌仕組み", lx + pad, cy_l, iw, C_PURPLE)
    t(slide,
      "テキストを数値ベクトルに変換し、コサイン類似度で「意味が近いもの」を取得。\n"
      "キーワードが違っても意味が近ければヒットする。",
      lx + pad, cy_l + Inches(0.26), iw, Inches(0.70),
      size=Pt(10), color=C_NAVY, align=PP_ALIGN.LEFT)

    # 検索対象
    section_label(slide, "▌検索対象（3種類のデータ）", lx + pad, cy_l + Inches(1.02), iw, C_PURPLE)
    src_h = Inches(0.34)
    for i, (lbl, sub) in enumerate([
        ("external.json", "市場・競合・規制データ"),
        ("internal.json", "技術・過去PJ・失敗事例"),
        ("persons.json",  "社内キーマン情報"),
    ]):
        by = cy_l + Inches(1.28) + i * (src_h + Inches(0.06))
        chip(slide, lbl, lx + pad, by, iw, src_h,
             fill=C_PURPLE, sub=sub)

    # 具体例
    section_label(slide, "▌具体例：テーマ「ビルエネルギー管理」で検索すると…",
                  lx + pad, cy_l + Inches(2.72), iw, C_PURPLE)
    eg_items = [
        "✓  BEMS市場の成長動向",
        "✓  過去のBEMS事業（失敗事例）",
        "✓  薄膜太陽電池技術（エネルギーで意味一致）",
    ]
    for i, txt in enumerate(eg_items):
        t(slide, txt,
          lx + pad + Inches(0.05),
          cy_l + Inches(2.96) + i * Inches(0.36),
          iw - Inches(0.05), Inches(0.34),
          size=Pt(10), color=C_NAVY, align=PP_ALIGN.LEFT)

    # 限界メモ
    rect(slide, lx + pad, cy_l + Inches(4.05), iw, Inches(0.30),
         fill=C_PALE_P, border=C_PURPLE, bw=Pt(1.0))
    t(slide, "⚠  「田村浩二」はテーマに出てこないためヒットしない → GraphRAG で補完",
      lx + pad + Inches(0.05), cy_l + Inches(4.05),
      iw - Inches(0.05), Inches(0.30),
      size=Pt(9), color=C_PURPLE, align=PP_ALIGN.LEFT, bold=True)

    # ════════════════════════════════
    # 中央：組み合わせ効果
    # ════════════════════════════════
    rect(slide, mx, panel_y, mw, panel_h, fill=C_NAVY)
    t(slide, "×", mx, panel_y + Inches(0.60), mw, Inches(0.80),
      size=Pt(48), color=C_WHITE, bold=True)
    t(slide, "組み合わせることで\n\n「関連度」\n　＋\n「関連性」\n\nを同時に取得",
      mx, panel_y + Inches(1.50), mw, Inches(2.40),
      size=Pt(10), color=C_TEAL, bold=True)
    harrow(slide,
           lx + lw + Inches(0.03), panel_y + panel_h / 2 - Inches(0.20),
           mx - Inches(0.03), color=C_PURPLE)
    harrow(slide,
           rx - Inches(0.03), panel_y + panel_h / 2 - Inches(0.20),
           mx + mw + Inches(0.03), color=C_RUST)
    t(slide, "↓\n抜け漏れのない\n検索を実現",
      mx, panel_y + Inches(4.00), mw, Inches(0.80),
      size=Pt(9), color=C_AMBER, bold=True)

    # ════════════════════════════════
    # 右パネル：GraphRAG（関連性）
    # ════════════════════════════════
    cy_r = panel(slide, "GraphRAG（NetworkX）── 関連性（つながり）で検索", rx, panel_y, rw, panel_h, C_RUST)

    section_label(slide, "▌仕組み", rx + pad, cy_r, iw, C_RUST)
    t(slide,
      "ChromaDB でヒットしたIDを起点に、知識グラフを1ステップ辿る。\n"
      "グラフ上でつながっている人物・技術・市場を自動で補完する。",
      rx + pad, cy_r + Inches(0.26), iw, Inches(0.70),
      size=Pt(10), color=C_NAVY, align=PP_ALIGN.LEFT)

    # グラフ構造の説明（簡易図）
    section_label(slide, "▌知識グラフの構造（ノードとエッジ）", rx + pad, cy_r + Inches(1.02), iw, C_RUST)
    node_types = [
        ("技術ノード",    "保有技術・コア資産"),
        ("人物ノード",    "社内キーマン・担当候補"),
        ("市場ノード",    "参入検討先の事業領域"),
        ("過去PJノード",  "実績・失敗事例"),
    ]
    nth = Inches(0.34)
    for i, (lbl, sub) in enumerate(node_types):
        by = cy_r + Inches(1.28) + i * (nth + Inches(0.06))
        chip(slide, lbl, rx + pad, by, iw, nth, fill=C_RUST, sub=sub)

    # 具体例
    section_label(slide, "▌具体例：GraphRAG が補完するもの",
                  rx + pad, cy_r + Inches(2.72), iw, C_RUST)
    t(slide,
      "ChromaDB ヒット：「BEMS事業（past_project）」\n"
      "      　　↓  グラフを1ステップ辿る\n"
      "✓  田村浩二（過去に担当・営業ネットワークあり）\n"
      "✓  薄膜太陽電池技術（応用できる）\n"
      "✓  BEMS市場（関連する）",
      rx + pad, cy_r + Inches(2.96), iw, Inches(1.50),
      size=Pt(10), color=C_NAVY, align=PP_ALIGN.LEFT)

    # conditions_now補完
    rect(slide, rx + pad, cy_r + Inches(4.05), iw, Inches(0.30),
         fill=C_PALE_R, border=C_RUST, bw=Pt(1.0))
    t(slide, "＋  internal.json を直接読込 → 失敗条件（conditions_now）を補完",
      rx + pad + Inches(0.05), cy_r + Inches(4.05),
      iw - Inches(0.05), Inches(0.30),
      size=Pt(9), color=C_RUST, align=PP_ALIGN.LEFT, bold=True)

    # ════════════════════════════════
    # 下段：比較表
    # ════════════════════════════════
    rect(slide, Inches(0.15), tbl_y, SW, tbl_h, fill=C_LIGHT, border=C_BORDER)

    cols = [
        ("検索方法",               Inches(2.50)),
        ("意味的に近い情報",       Inches(2.50)),
        ("構造的なつながり（関係者）", Inches(3.00)),
        ("失敗条件の解消確認",     Inches(2.70)),
        ("→ 抜け漏れリスク",      Inches(2.03)),
    ]
    hdr_h_tbl = Inches(0.28)
    row_h = (tbl_h - hdr_h_tbl - Inches(0.10)) / 3
    tx = Inches(0.25)
    for cname, cw in cols:
        rect(slide, tx, tbl_y, cw, hdr_h_tbl, fill=C_NAVY)
        t(slide, cname, tx, tbl_y, cw, hdr_h_tbl,
          size=Pt(9), color=C_WHITE, bold=True)
        tx += cw + Inches(0.03)

    rows = [
        ("キーワード検索",      "×", "×", "×", "高い"),
        ("ChromaDB のみ",       "○", "×", "×", "中程度"),
        ("ChromaDB ＋ GraphRAG","○", "○", "○", "低い"),
    ]
    row_colors = [C_NAVY, C_NAVY, C_TEAL]
    marks_color = [
        (RGBColor(0xEF, 0x44, 0x44), RGBColor(0xEF, 0x44, 0x44), RGBColor(0xEF, 0x44, 0x44), RGBColor(0xEF, 0x44, 0x44)),
        (C_GREEN, RGBColor(0xEF, 0x44, 0x44), RGBColor(0xEF, 0x44, 0x44), C_AMBER),
        (C_GREEN, C_GREEN, C_GREEN, C_GREEN),
    ]

    for ri, (row, mc) in enumerate(zip(rows, marks_color)):
        ry = tbl_y + hdr_h_tbl + Inches(0.05) + ri * (row_h + Inches(0.02))
        tx = Inches(0.25)
        for ci, (val, (cname, cw)) in enumerate(zip(row, cols)):
            bg = row_colors[ri] if ci == 0 else (C_PALE_P if ri == 1 and ci > 0 else (C_NAVY if ri == 2 else C_LIGHT))
            rect(slide, tx, ry, cw, row_h, fill=bg)
            fc = C_WHITE if bg in (C_NAVY, C_TEAL) else mc[ci - 1] if ci > 0 else C_WHITE
            if ci == 0:
                t(slide, val, tx, ry, cw, row_h, size=Pt(10), color=C_WHITE, bold=True)
            elif ri == 2:
                t(slide, val, tx, ry, cw, row_h, size=Pt(13), color=C_GREEN, bold=True)
            elif val == "×":
                t(slide, val, tx, ry, cw, row_h, size=Pt(13),
                  color=RGBColor(0xEF, 0x44, 0x44), bold=True)
            else:
                t(slide, val, tx, ry, cw, row_h, size=Pt(13), color=C_GREEN, bold=True)
            tx += cw + Inches(0.03)

    ftr(slide, "関連度（ChromaDB）＋ 関連性（GraphRAG）の組み合わせで初めて、抜け漏れのない情報収集が可能になる")


# ══════════════════════════════════════════════════════════════════════
# Slide 2 : AI分析への接続と可視化
# ══════════════════════════════════════════════════════════════════════

def build_slide2(prs):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG)

    hdr(slide, "収集した情報をAIが分析し、必要な情報だけを可視化する")

    cy = HDR_Y + HDR_H + Inches(0.14)
    bh = Inches(0.38)
    bg = Inches(0.12)

    # ── 4ゾーンのX座標 ──
    z1x, z1w = Inches(0.15), Inches(2.80)    # 3軸コンテキスト
    z2x, z2w = Inches(3.20), Inches(3.30)    # Stage1
    z3x, z3w = Inches(6.75), Inches(2.50)    # GO/NO + Stage2
    z4x, z4w = Inches(9.50), Inches(3.68)    # 出力・可視化

    zone_h = FTR_Y - cy - Inches(0.08)

    # ── Zone1：3軸コンテキスト ──
    rect(slide, z1x, cy, z1w, zone_h, fill=C_LIGHT, border=C_BORDER)
    rect(slide, z1x, cy, z1w, Inches(0.32), fill=C_TEAL)
    t(slide, "3軸コンテキスト", z1x + Inches(0.08), cy, z1w, Inches(0.32),
      size=Pt(11), color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)
    t(slide, "ChromaDB ＋ GraphRAG\nの収集結果を整形",
      z1x + Inches(0.08), cy + Inches(0.36), z1w - Inches(0.16), Inches(0.55),
      size=Pt(9), color=C_GRAY, align=PP_ALIGN.LEFT)

    ctx_items = [
        ("外部コンテキスト", "市場・競合・規制"),
        ("社内コンテキスト", "技術・過去PJ\n失敗条件（conditions_now）"),
        ("組織コンテキスト", "キーマン情報\n（グラフ補完分含む）"),
    ]
    ctx_bh = Inches(1.10)
    for i, (lbl, sub) in enumerate(ctx_items):
        by = cy + Inches(0.98) + i * (ctx_bh + Inches(0.12))
        chip(slide, lbl, z1x + Inches(0.10), by,
             z1w - Inches(0.20), ctx_bh, fill=C_TEAL, sub=sub)

    # ── Zone1 → Zone2 矢印 ──
    arr_y = cy + zone_h / 2
    harrow(slide, z1x + z1w + Inches(0.04), arr_y,
           z2x - Inches(0.04), color=C_TEAL)

    # ── Zone2：Stage1（3並列評価）──
    rect(slide, z2x, cy, z2w, zone_h, fill=C_LIGHT, border=C_BORDER)
    rect(slide, z2x, cy, z2w, Inches(0.32), fill=C_NAVY)
    t(slide, "Stage1　3軸を個別評価（3並列）",
      z2x + Inches(0.08), cy, z2w, Inches(0.32),
      size=Pt(11), color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)
    t(slide, "3軸を独立したGPT呼び出しで同時処理",
      z2x + Inches(0.08), cy + Inches(0.36), z2w - Inches(0.16), Inches(0.40),
      size=Pt(9), color=C_GRAY, align=PP_ALIGN.LEFT)

    s1_items = [
        ("外部環境を評価", "市場・規制・競合の観点"),
        ("社内適合を評価", "技術・失敗条件の観点"),
        ("組織体制を評価", "キーマン・体制の観点"),
    ]
    s1_bh = Inches(1.10)
    for i, (lbl, sub) in enumerate(s1_items):
        by = cy + Inches(0.84) + i * (s1_bh + Inches(0.12))
        chip(slide, lbl, z2x + Inches(0.10), by,
             z2w - Inches(0.20), s1_bh, fill=C_NAVY, sub=sub)
        # スコアバッジ
        rect(slide, z2x + z2w - Inches(0.58), by + Inches(0.06),
             Inches(0.44), Inches(0.32), fill=C_AMBER)
        t(slide, "◎○△×",
          z2x + z2w - Inches(0.58), by + Inches(0.06),
          Inches(0.44), Inches(0.32),
          size=Pt(8), color=C_WHITE, bold=True)

    # ── Zone2 → Zone3 矢印 ──
    harrow(slide, z2x + z2w + Inches(0.04), arr_y,
           z3x - Inches(0.04), color=C_NAVY, lbl="スコア")

    # ── Zone3：GO/NO判定 ＋ Stage2 ──
    rect(slide, z3x, cy, z3w, zone_h, fill=C_LIGHT, border=C_BORDER)
    rect(slide, z3x, cy, z3w, Inches(0.32), fill=C_NAVY)
    t(slide, "判定 ＋ Stage2",
      z3x + Inches(0.08), cy, z3w, Inches(0.32),
      size=Pt(11), color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)

    # GO/NO
    gono_h = Inches(1.10)
    rect(slide, z3x + Inches(0.10), cy + Inches(0.44),
         z3w - Inches(0.20), gono_h, fill=C_GREEN)
    t(slide, "GO / NO 判定",
      z3x + Inches(0.10), cy + Inches(0.44),
      z3w - Inches(0.20), Inches(0.36),
      size=Pt(11), color=C_WHITE, bold=True)
    t(slide, "ルールベース（主観なし）\n社内適合×→NO\n社内△→条件付きGO",
      z3x + Inches(0.10), cy + Inches(0.44) + Inches(0.34),
      z3w - Inches(0.20), Inches(0.72),
      size=Pt(9), color=C_WHITE, align=PP_ALIGN.LEFT)

    # Stage2 Tier1/Tier2
    t2_y = cy + Inches(0.44) + gono_h + Inches(0.18)
    t(slide, "Stage2（Tier1 ＋ Tier2 並列）",
      z3x + Inches(0.08), t2_y - Inches(0.22),
      z3w - Inches(0.16), Inches(0.22),
      size=Pt(9), color=C_TEAL, bold=True, align=PP_ALIGN.LEFT)
    tier_h = Inches(0.95)
    chip(slide, "Tier1",
         z3x + Inches(0.10), t2_y,
         z3w - Inches(0.20), tier_h, fill=C_NAVY,
         sub="事業案3つ ＋ 承認者サマリー")
    chip(slide, "Tier2",
         z3x + Inches(0.10), t2_y + tier_h + Inches(0.12),
         z3w - Inches(0.20), tier_h, fill=C_NAVY,
         sub="3C分析\nCustomer / Competitor / Company")

    # ── Zone3 → Zone4 矢印 ──
    harrow(slide, z3x + z3w + Inches(0.04), arr_y,
           z4x - Inches(0.04), color=C_NAVY, lbl="生成")

    # ── Zone4：出力・可視化 ──
    rect(slide, z4x, cy, z4w, zone_h, fill=C_LIGHT, border=C_BORDER)
    rect(slide, z4x, cy, z4w, Inches(0.32), fill=C_TEAL)
    t(slide, "出力・可視化",
      z4x + Inches(0.08), cy, z4w, Inches(0.32),
      size=Pt(11), color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)
    t(slide, "必要な情報だけを可視化",
      z4x + Inches(0.08), cy + Inches(0.36), z4w - Inches(0.16), Inches(0.30),
      size=Pt(9), color=C_GRAY, align=PP_ALIGN.LEFT)

    out_items = [
        ("承認者サマリー",     "GO/NOの理由を1段落に凝縮\n意思決定者が読むべき情報のみ"),
        ("3軸評価パネル",      "◎○△×スコア ＋ 評価根拠\n判断の経緯を後から追える"),
        ("事業提案タブ",       "事業案3つ ＋ 担当者アクション案"),
        ("PyVisグラフ",        "参照ノードを黄色でハイライト\n「なぜこの人物・技術か」が一目でわかる"),
    ]
    out_bh = Inches(0.88)
    for i, (lbl, sub) in enumerate(out_items):
        by = cy + Inches(0.74) + i * (out_bh + Inches(0.10))
        chip(slide, lbl, z4x + Inches(0.10), by,
             z4w - Inches(0.20), out_bh, fill=C_TEAL, sub=sub)

    ftr(slide, "Stage1（3並列）→ GO/NO判定（ルールベース）→ Stage2（2並列）の設計により、高速かつ一貫した分析を実現")


# ══════════════════════════════════════════════════════════════════════
# メイン
# ══════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    build_slide1(prs)
    build_slide2(prs)
    prs.save(OUTPUT_PATH)
    print(f"保存完了: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
