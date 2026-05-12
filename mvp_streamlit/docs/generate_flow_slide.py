"""
アーキテクチャフロー スライド生成（縦型・Step2/Step3詳細表示版）
実行: python docs/generate_flow_slide.py
出力: docs/flow_slide.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUTPUT_PATH = "docs/flow_slide.pptx"

# ── カラーパレット ──
C_BG      = RGBColor(0xFF, 0xFF, 0xFF)  # 白背景
C_NAVY    = RGBColor(0x1B, 0x2E, 0x45)  # 濃紺（ステップヘッダー）
C_TEAL    = RGBColor(0x00, 0x99, 0xCC)  # ティール（ラベル）
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_ORANGE  = RGBColor(0xFF, 0x88, 0x22)  # input矢印
C_GPT     = RGBColor(0x10, 0xA3, 0x7F)  # OpenAI
C_PURPLE  = RGBColor(0x7C, 0x3A, 0xED)  # ChromaDB
C_RUST    = RGBColor(0xC0, 0x50, 0x10)  # NetworkX / グラフ
C_RED     = RGBColor(0xB9, 0x1C, 0x1C)  # 失敗条件補完
C_GREEN   = RGBColor(0x16, 0xA3, 0x4A)  # GO
C_CRIMSON = RGBColor(0xDC, 0x26, 0x26)  # NO
C_GRAY    = RGBColor(0x9C, 0xA3, 0xAF)  # サブテキスト
C_STEP_BG = RGBColor(0xF8, 0xFA, 0xFC)  # ステップ背景
C_BORDER  = RGBColor(0xD1, 0xD5, 0xDB)  # ステップ枠

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── ステップボックス共通設定 ──
SX   = Inches(0.15)                   # ステップ左端
SW   = SLIDE_W - Inches(0.30)         # ステップ幅
STH  = Inches(0.30)                   # ステップタイトル帯の高さ
GAP  = Inches(0.10)                   # ステップ間ギャップ
HDR_H = Inches(0.40)
HDR_Y = Inches(0.10)
FTR_H = Inches(0.40)
FTR_Y = SLIDE_H - FTR_H

# ── ステップ高さ ──
S1_H = Inches(0.68)
S2_H = Inches(1.40)
S3_H = Inches(1.80)
S4_H = Inches(1.42)
S5_H = Inches(0.66)

S1_Y = HDR_Y + HDR_H + Inches(0.08)
S2_Y = S1_Y + S1_H + GAP
S3_Y = S2_Y + S2_H + GAP
S4_Y = S3_Y + S3_H + GAP
S5_Y = S4_Y + S4_H + GAP


# ── ユーティリティ ──────────────────────────────────

def rect(slide, x, y, w, h, fill=None, border=None, bw=Pt(1.2), dashed=False):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
        if dashed:
            ln = s.line._ln
            for c in ln.findall(qn('a:prstDash')): ln.remove(c)
            pd = etree.SubElement(ln, qn('a:prstDash')); pd.set('val', 'dash')
    else:
        s.line.fill.background()
    return s


def oval(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(9, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def t(slide, text, x, y, w, h, size=Pt(11), color=None,
      bold=False, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold
    r.font.color.rgb = color or C_NAVY
    return tb


def chip(slide, label, x, y, w, h, fill=None, sub=None):
    """濃紺ボックス（オプションでサブテキスト付き）"""
    bg = fill or C_NAVY
    rect(slide, x, y, w, h, fill=bg)
    if sub:
        t(slide, label, x, y, w, h * 0.52,
          size=Pt(10), color=C_WHITE, bold=True)
        t(slide, sub, x, y + h * 0.50, w, h * 0.44,
          size=Pt(8), color=RGBColor(0xCC, 0xE5, 0xFF))
    else:
        t(slide, label, x, y, w, h, size=Pt(10), color=C_WHITE, bold=True)


def arrowhead(conn):
    ln = conn.line._ln
    el = ln.find(qn('a:tailEnd'))
    if el is None:
        el = etree.SubElement(ln, qn('a:tailEnd'))
    el.set('type', 'arrow'); el.set('w', 'med'); el.set('len', 'med')


def harrow(slide, x1, y, x2, color=None, lbl=None, lbl_above=True):
    c = color or C_NAVY
    co = slide.shapes.add_connector(1, x1, y, x2, y)
    co.line.color.rgb = c; co.line.width = Pt(1.5)
    arrowhead(co)
    if lbl:
        lx, lw = min(x1, x2), abs(x2 - x1)
        ly = y - Inches(0.20) if lbl_above else y + Inches(0.03)
        t(slide, lbl, lx, ly, lw, Inches(0.18), size=Pt(8), color=c, bold=True)


def varrow(slide, x, y1, y2, color=None, lbl=None):
    c = color or C_NAVY
    co = slide.shapes.add_connector(1, x, y1, x, y2)
    co.line.color.rgb = c; co.line.width = Pt(1.5)
    arrowhead(co)
    if lbl:
        t(slide, lbl, x + Inches(0.05), (y1 + y2) / 2 - Inches(0.10),
          Inches(1.0), Inches(0.20), size=Pt(8), color=c, bold=True,
          align=PP_ALIGN.LEFT)


def step_frame(slide, num, title, y, h, color=C_NAVY):
    """ステップ外枠＋タイトル帯を描画する"""
    # 外枠
    rect(slide, SX, y, SW, h, fill=C_STEP_BG, border=C_BORDER)
    # タイトル帯（濃紺）
    rect(slide, SX, y, SW, STH, fill=color)
    # 番号
    t(slide, num, SX + Inches(0.06), y, Inches(0.48), STH,
      size=Pt(14), color=C_WHITE, bold=True)
    # タイトル
    t(slide, title, SX + Inches(0.55), y, SW - Inches(0.60), STH,
      size=Pt(12), color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)


# ── コンテンツ領域の左端・幅 ──
CX = SX + Inches(0.20)                  # コンテンツ左端
CW = SW - Inches(0.40)                  # コンテンツ幅


def content_top(step_y):
    """ステップ内のコンテンツ開始Y"""
    return step_y + STH + Inches(0.10)


# ═══════════════════════════════════════════════════
# メイン
# ═══════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG)

    # ── ヘッダー（ワンメッセージ）──
    rect(slide, Inches(0.15), HDR_Y, SW, HDR_H, fill=C_NAVY)
    t(slide, "テーマを入力するだけで、数日かかっていた情報収集・整理・判断を数十秒で完了する",
      Inches(0.15), HDR_Y, SW, HDR_H,
      size=Pt(14), color=C_WHITE, bold=True)

    # ════════════════════════════════════════════════
    # Step 1：テーマ入力
    # ════════════════════════════════════════════════
    step_frame(slide, "①", "テーマ入力", S1_Y, S1_H)
    cy1 = content_top(S1_Y)
    ih1 = S1_H - STH - Inches(0.18)
    bw1 = Inches(3.6)
    for i, lbl in enumerate([
        "ターゲット市場 / 想定顧客",
        "活用したい自社アセット・技術",
        "事業アイデアの詳細（＋PDF添付）",
    ]):
        bx = CX + i * (bw1 + Inches(0.18))
        chip(slide, lbl, bx, cy1, bw1, ih1)

    # ════════════════════════════════════════════════
    # ステップ間矢印
    # ════════════════════════════════════════════════
    varrow(slide, SLIDE_W / 2, S1_Y + S1_H + Inches(0.01), S2_Y - Inches(0.01))

    # ════════════════════════════════════════════════
    # Step 2：データ検索（ChromaDB）
    # ════════════════════════════════════════════════
    step_frame(slide, "②", "データ検索 ── ChromaDB による意味検索", S2_Y, S2_H, C_PURPLE)
    cy2 = content_top(S2_Y)
    ih2 = S2_H - STH - Inches(0.20)
    src_h = (ih2 - Inches(0.10)) / 3   # 3つのソースボックスの高さ

    # ── 左：3つのデータソース ──
    src_x = CX
    src_w = Inches(2.80)
    src_labels = [
        ("external.json", "市場・競合・規制データ"),
        ("internal.json", "技術・過去PJ・失敗事例"),
        ("persons.json",  "社内キーマン情報"),
    ]
    for i, (title, sub) in enumerate(src_labels):
        by = cy2 + i * (src_h + Inches(0.05))
        chip(slide, title, src_x, by, src_w, src_h - Inches(0.02), sub=sub)

    # ── 矢印 → ChromaDB ──
    db_x = src_x + src_w + Inches(0.50)
    db_w = Inches(3.20)
    db_mid_y = cy2 + ih2 / 2
    harrow(slide, src_x + src_w + Inches(0.05), db_mid_y,
           db_x - Inches(0.05), color=C_PURPLE)

    # ── 中央：ChromaDB ──
    chip(slide, "ChromaDB", db_x, cy2, db_w, ih2,
         fill=C_PURPLE,
         sub="意味の近さ（コサイン類似度）で\nキーワード検索より広く・正確に取得")

    # ── 矢印 → 結果 ──
    res_x = db_x + db_w + Inches(0.50)
    res_w = CX + CW - res_x
    harrow(slide, db_x + db_w + Inches(0.05), db_mid_y,
           res_x - Inches(0.05), color=C_PURPLE, lbl="上位5件を取得")

    # ── 右：検索結果 ──
    chip(slide, "検索結果（上位5件）", res_x, cy2, res_w, ih2,
         sub="id / 内容テキスト\n類似度スコア / sourceタグ")

    # ════════════════════════════════════════════════
    # ステップ間矢印
    # ════════════════════════════════════════════════
    varrow(slide, SLIDE_W / 2, S2_Y + S2_H + Inches(0.01), S3_Y - Inches(0.01))

    # ════════════════════════════════════════════════
    # Step 3：コンテキスト構築（2つの補完処理）
    # ════════════════════════════════════════════════
    step_frame(slide, "③", "コンテキスト構築 ── 2つの補完処理で3軸テキストを生成", S3_Y, S3_H, C_RUST)
    cy3 = content_top(S3_Y)
    ih3 = S3_H - STH - Inches(0.20)

    # ── 左：入力（Step2の結果） ──
    inp_w = Inches(2.20)
    inp_x = CX
    chip(slide, "Step②の\n検索結果（5件）",
         inp_x, cy3, inp_w, ih3, fill=C_PURPLE)

    # ── 中央：2つの補完処理（上下） ──
    proc_x = inp_x + inp_w + Inches(0.55)
    proc_w = Inches(3.80)
    proc_h = (ih3 - Inches(0.18)) / 2

    # 上：グラフ検索
    gx_y = cy3
    chip(slide, "補完① グラフ検索（NetworkX）",
         proc_x, gx_y, proc_w, proc_h,
         fill=C_RUST,
         sub="ヒットIDを起点にノードを1ステップ辿る\n→ 直接ヒットしなかった関係者・技術を補完")

    # 下：失敗条件補完
    enrich_y = gx_y + proc_h + Inches(0.18)
    chip(slide, "補完② 失敗条件チェック（internal.json 直接読込）",
         proc_x, enrich_y, proc_w, proc_h,
         fill=C_RED,
         sub="conditions_now を取得\n→ 過去の失敗条件が現在解消されたかを確認")

    # 矢印：入力 → 各補完処理
    in_mid_y = cy3 + ih3 / 2
    harrow(slide, inp_x + inp_w + Inches(0.05), gx_y + proc_h / 2,
           proc_x - Inches(0.05), color=C_RUST)
    harrow(slide, inp_x + inp_w + Inches(0.05), enrich_y + proc_h / 2,
           proc_x - Inches(0.05), color=C_RED)

    # ── 右：3軸コンテキスト出力 ──
    out_x = proc_x + proc_w + Inches(0.55)
    out_w = CX + CW - out_x
    out_labels = [
        ("外部コンテキスト", "市場・競合・規制"),
        ("社内コンテキスト", "技術・過去PJ・失敗条件"),
        ("組織コンテキスト", "キーマン情報"),
    ]
    out_bh = (ih3 - Inches(0.10)) / 3
    for i, (title, sub) in enumerate(out_labels):
        by = cy3 + i * (out_bh + Inches(0.05))
        chip(slide, title, out_x, by, out_w, out_bh - Inches(0.02),
             fill=C_TEAL, sub=sub)

    # 矢印：補完処理 → 出力
    harrow(slide, proc_x + proc_w + Inches(0.05), gx_y + proc_h / 2,
           out_x - Inches(0.05), color=C_TEAL, lbl="3軸に整形")
    harrow(slide, proc_x + proc_w + Inches(0.05), enrich_y + proc_h / 2,
           out_x - Inches(0.05), color=C_TEAL)

    # ════════════════════════════════════════════════
    # ステップ間矢印
    # ════════════════════════════════════════════════
    varrow(slide, SLIDE_W / 2, S3_Y + S3_H + Inches(0.01), S4_Y - Inches(0.01))

    # ════════════════════════════════════════════════
    # Step 4：AI分析（Stage1 → GO/NO → Stage2）
    # ════════════════════════════════════════════════
    step_frame(slide, "④", "AI分析 ── 2段階・並列処理", S4_Y, S4_H)
    cy4 = content_top(S4_Y)
    ih4 = S4_H - STH - Inches(0.20)

    zone_w = (CW - Inches(1.0)) / 3
    z1_x = CX
    z2_x = z1_x + zone_w + Inches(0.50)
    z3_x = z2_x + zone_w + Inches(0.50)
    small_h = (ih4 - Inches(0.10)) / 3

    # Stage1（3並列）
    for i, lbl in enumerate(["外部環境を評価", "社内適合を評価", "組織体制を評価"]):
        by = cy4 + i * (small_h + Inches(0.05))
        chip(slide, lbl, z1_x, by, zone_w, small_h - Inches(0.02))
    t(slide, "Stage1（3並列）",
      z1_x, cy4 - Inches(0.02), zone_w, Inches(0.20),
      size=Pt(9), color=C_TEAL, bold=True)

    # 矢印 → GO/NO
    harrow(slide, z1_x + zone_w + Inches(0.05), cy4 + ih4 / 2,
           z2_x - Inches(0.05), lbl="◎○△× スコア")

    # GO/NO判定
    chip(slide, "GO / NO 判定", z2_x, cy4, zone_w, ih4 * 0.55,
         fill=C_GREEN,
         sub="社内適合×→NO\n社内△→条件付きGO\nそれ以外→GO")
    t(slide, "ルールベース（主観なし）",
      z2_x, cy4 - Inches(0.02), zone_w, Inches(0.20),
      size=Pt(9), color=C_GREEN, bold=True)

    # 矢印 → Stage2
    harrow(slide, z2_x + zone_w + Inches(0.05), cy4 + ih4 / 2,
           z3_x - Inches(0.05), lbl="GO/NO判定を付与")

    # Stage2（2並列）
    chip(slide, "Tier1：事業案3つ ＋ 承認者サマリー",
         z3_x, cy4, zone_w, small_h * 1.4 - Inches(0.04))
    chip(slide, "Tier2：3C分析",
         z3_x, cy4 + small_h * 1.4 + Inches(0.08),
         zone_w, small_h * 1.4 - Inches(0.04),
         sub="Customer / Competitor / Company")
    t(slide, "Stage2（2並列）",
      z3_x, cy4 - Inches(0.02), zone_w, Inches(0.20),
      size=Pt(9), color=C_TEAL, bold=True)

    # ════════════════════════════════════════════════
    # ステップ間矢印
    # ════════════════════════════════════════════════
    varrow(slide, SLIDE_W / 2, S4_Y + S4_H + Inches(0.01), S5_Y - Inches(0.01))

    # ════════════════════════════════════════════════
    # Step 5：画面表示
    # ════════════════════════════════════════════════
    step_frame(slide, "⑤", "画面表示（Streamlit）", S5_Y, S5_H)
    cy5 = content_top(S5_Y)
    ih5 = S5_H - STH - Inches(0.18)
    disp = [
        "承認者サマリー\n（GO/NO ＋ 要約）",
        "3軸評価\n（◎○△× ＋ 根拠）",
        "事業提案タブ\n（3案 ＋ アクション）",
        "3C分析タブ",
        "知識グラフ\n（PyVis）",
    ]
    dw = (CW - Inches(0.20) * (len(disp) - 1)) / len(disp)
    for i, lbl in enumerate(disp):
        bx = CX + i * (dw + Inches(0.20))
        chip(slide, lbl, bx, cy5, dw, ih5)

    # ── フッター ──
    rect(slide, Inches(0.15), FTR_Y, SW, FTR_H, fill=C_NAVY)
    t(slide, "ChromaDB（意味検索）＋ NetworkX（グラフ検索）＋ GPT-4o-mini（AI分析）を組み合わせた新規事業判断支援システム",
      Inches(0.15), FTR_Y, SW, FTR_H,
      size=Pt(11), color=C_WHITE)

    prs.save(OUTPUT_PATH)
    print(f"保存完了: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
